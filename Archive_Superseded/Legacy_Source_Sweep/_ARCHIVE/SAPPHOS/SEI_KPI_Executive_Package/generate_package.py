import os
from docx import Document
from docx.shared import Inches
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches as PptInches

# Paths
base_folder = r'C:\Users\mespindola\Desktop\SEI_KPI_Executive_Package'
charts_folder = os.path.join(base_folder, 'charts')
docx_folder = os.path.join(base_folder, 'docx')
pptx_path = os.path.join(base_folder, 'SEI_Executive_KPI_Report.pptx')

# Create folders if missing
for path in [charts_folder, docx_folder]:
    if not os.path.exists(path):
        os.makedirs(path)

# ---------- Generate sample charts ----------
# KPI Bar Chart
employees = ['Alice', 'Bob', 'Charlie', 'David']
chargeability = [85, 78, 90, 70]
plt.figure(figsize=(6,4))
plt.bar(employees, chargeability, color='skyblue')
plt.title('Chargeability by Employee')
plt.ylabel('Chargeability (%)')
kpi_chart_path = os.path.join(charts_folder, 'kpi_bar_chart.png')
plt.savefig(kpi_chart_path)
plt.close()

# Radar Chart
labels = ['Work Quality','Client Satisfaction','Project Delivery','Professional Dev','Compliance']
stats = [80, 70, 75, 60, 90]
angles = np.linspace(0, 2 * np.pi, len(labels), endpoint=False).tolist()
stats += stats[:1]
angles += angles[:1]
fig = plt.figure(figsize=(5,5))
ax = fig.add_subplot(111, polar=True)
ax.plot(angles, stats, 'o-', linewidth=2)
ax.fill(angles, stats, alpha=0.25)
ax.set_thetagrids(np.degrees(angles), labels)
radar_chart_path = os.path.join(charts_folder, 'radar_chart.png')
plt.savefig(radar_chart_path)
plt.close()

# Equity Heatmap
roles = ['Tech','PM','Leader']
gender = ['M','F']
data = pd.DataFrame([[80,70],[85,75],[60,65]], index=roles, columns=gender)
plt.figure(figsize=(4,3))
plt.imshow(data, cmap='coolwarm', aspect='auto')
plt.colorbar(label='Chargeability %')
plt.xticks(range(len(gender)), gender)
plt.yticks(range(len(roles)), roles)
equity_heatmap_path = os.path.join(charts_folder, 'equity_heatmap.png')
plt.savefig(equity_heatmap_path)
plt.close()

# ---------- Generate sample DOCX documents ----------
docs = {
    "Policies.docx": ["Protected Leave Policy","Company Leave Policy","KPI Calculation Policy"],
    "KPI_Workflow.docx": ["Step 1: Time Entry","Step 2: HR Leave Adjustment","Step 3: KPI Calculation","Step 4: Executive Review"],
    "KPI_Definitions.docx": ["Chargeability","Protected Leave","Non-Protected Leave","Core Competencies","KPI","Job Role & Classification","Employee Performance"]
}

for doc_name, sections in docs.items():
    doc_path = os.path.join(docx_folder, doc_name)
    doc = Document()
    doc.add_heading(doc_name.replace('.docx',''),0)
    for sec in sections:
        doc.add_heading(sec,level=1)
        doc.add_paragraph("High-level description for " + sec)
    doc.save(doc_path)

# ---------- Generate PowerPoint ----------
prs = Presentation()
# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "SEI Executive KPI Report"
slide.placeholders[1].text = "High-Level Overview"

# Slide 2: KPI Bar Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Chargeability by Employee"
slide.shapes.add_picture(kpi_chart_path, PptInches(1), PptInches(1), width=PptInches(6))

# Slide 3: Radar Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "KPI Balance Radar"
slide.shapes.add_picture(radar_chart_path, PptInches(1), PptInches(1), width=PptInches(6))

# Slide 4: Equity Heatmap
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Equity Heatmap"
slide.shapes.add_picture(equity_heatmap_path, PptInches(1), PptInches(1), width=PptInches(6))

prs.save(pptx_path)
print(f"✅ PowerPoint generated at: {pptx_path}")
print(f"✅ DOCX documents generated at: {docx_folder}")
