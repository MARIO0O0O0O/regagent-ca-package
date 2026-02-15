import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import numpy as np

# Paths
report_path = r'C:\Users\mespindola\Desktop\SEI_KPI_Executive_Presentation\SEI_Executive_KPI_Report.pptx'
chart_folder = r'C:\Users\mespindola\Desktop\SEI_KPI_Executive_Presentation\charts'

if not os.path.exists(chart_folder):
    os.makedirs(chart_folder)

# Create Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ---------- Helper Functions ----------
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide

def add_image_slide(prs, title, img_path):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(11))
    return slide

# ---------- 1. Title Slide ----------
add_title_slide(prs,
    "SEI Executive KPI Report – High-Level Overview",
    "Prepared for Executive Leadership\nPrepared by SEI HR & Performance Analytics"
)

# ---------- 2. Industry Benchmark Slide ----------
add_bullet_slide(prs,
    "1. Industry Benchmark Intelligence",
    [
        "Technical / Junior Consultants: 75–85% utilization (non-billable dev time allowed)",
        "Mid-Level Consultants / PMs: 75–90% (higher client responsibility, some internal work)",
        "Senior Leaders / Partners: 40–75% (strategic & business development focus)",
        "Firm Averages: 65–80% (balance revenue + growth)",
        "Role-specific ranges increase transparency and employee buy-in",
        "Goldilocks Zone (2025 benchmark): 70–80% chargeability"
    ]
)

# ---------- 3. KPI Bar Chart ----------
employees = ['Alice','Bob','Charlie','David']
chargeability = [85,78,90,70]
kpi_chart_path = os.path.join(chart_folder, 'kpi_bar_chart.png')
plt.figure(figsize=(10,5))
bars = plt.bar(employees, chargeability, color='skyblue', edgecolor='black')
plt.ylim(0,100)
plt.title("High-Level Chargeability by Employee", fontsize=16)
plt.ylabel("Chargeability (%)", fontsize=12)
plt.grid(axis='y', linestyle='--', alpha=0.7)
for bar, val in zip(bars, chargeability):
    plt.text(bar.get_x() + bar.get_width()/2, val+1, f'{val}%', ha='center', fontsize=10)
plt.tight_layout()
plt.savefig(kpi_chart_path)
plt.close()
add_image_slide(prs, "2. KPI Chargeability Overview", kpi_chart_path)

# ---------- 4. Radar Chart for KPI Balance ----------
kpi_categories = ['Work Quality','Client Satisfaction','Project Delivery','Professional Development','Compliance']
kpi_scores = [90,85,80,70,75]
angles = np.linspace(0, 2 * np.pi, len(kpi_categories), endpoint=False).tolist()
scores = kpi_scores + [kpi_scores[0]]
angles += [angles[0]]

radar_path = os.path.join(chart_folder,'kpi_radar.png')
fig, ax = plt.subplots(figsize=(6,6), subplot_kw=dict(polar=True))
ax.plot(angles, scores, color='teal', linewidth=2, linestyle='solid')
ax.fill(angles, scores, color='teal', alpha=0.25)
ax.set_xticks(angles[:-1])
ax.set_xticklabels(kpi_categories)
ax.set_yticks([20,40,60,80,100])
ax.set_yticklabels(['20','40','60','80','100'])
plt.title("KPI Balance Radar Chart", fontsize=14)
plt.savefig(radar_path)
plt.close()
add_image_slide(prs, "3. KPI Balance Radar Chart", radar_path)

# ---------- 5. Equity Heatmap ----------
roles = ['Technical','PM','Leadership']
genders = ['Male','Female']
data = np.array([[80,75],[70,72],[60,65]])  # sample chargeability % by gender/role
heatmap_path = os.path.join(chart_folder,'equity_heatmap.png')
fig, ax = plt.subplots()
c = ax.imshow(data, cmap='YlOrRd', vmin=0, vmax=100)
ax.set_xticks(np.arange(len(genders)))
ax.set_yticks(np.arange(len(roles)))
ax.set_xticklabels(genders)
ax.set_yticklabels(roles)
for i in range(len(roles)):
    for j in range(len(genders)):
        ax.text(j, i, f"{data[i,j]}%", ha='center', va='center', color='black', fontsize=10)
plt.title("Equity Heatmap: Chargeability by Role & Gender")
plt.colorbar(c)
plt.savefig(heatmap_path)
plt.close()
add_image_slide(prs, "4. Equity Heatmap", heatmap_path)

# ---------- 6. Gantt-like Timeline ----------
tasks = ['Policy Design','Chart Generation','Manager Training','Employee Communication','Quarterly Audit']
start = [0,2,4,6,8]
dur = [2,2,2,2,2]
gantt_path = os.path.join(chart_folder,'gantt_timeline.png')
plt.figure(figsize=(10,3))
for i, task in enumerate(tasks):
    plt.barh(task, dur[i], left=start[i], color='mediumseagreen', edgecolor='black')
plt.xlabel('Months')
plt.title("High-Level KPI Implementation Timeline")
plt.tight_layout()
plt.savefig(gantt_path)
plt.close()
add_image_slide(prs, "5. KPI Implementation Timeline", gantt_path)

# ---------- 7. Executive Summary ----------
add_bullet_slide(prs,
    "6. Executive Summary",
    [
        "✅ High-level KPI formula ensures fairness by adjusting for protected & company leaves",
        "✅ Role-specific targets balance profitability & employee expectations",
        "✅ Broader KPIs complement chargeability (e.g., project margin, realization rate)",
        "✅ Equity & compliance audits mitigate risk and support defensible decision-making",
        "✅ Time tracking and dashboards ensure transparency & actionable insights"
    ]
)

# Save Presentation
prs.save(report_path)
print(f"✅ PowerPoint generated: {report_path}")
