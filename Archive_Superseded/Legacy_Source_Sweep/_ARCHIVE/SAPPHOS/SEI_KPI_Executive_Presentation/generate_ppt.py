import os

# Self-healing imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print('⚠️ python-pptx not installed. PowerPoint generation skipped.')
    exit()

try:
    import matplotlib.pyplot as plt
except ImportError:
    print('⚠️ matplotlib not installed. Charts will be placeholders.')
    plt = None

ppt_path = r'C:\Users\mespindola\Desktop\SEI_KPI_Executive_Presentation\SEI_Executive_KPI_Report.pptx'
chart_folder = r'C:\Users\mespindola\Desktop\SEI_KPI_Executive_Presentation\charts'

# Create presentation
prs = Presentation()

# Function to add slide
def add_slide(title, content=None, table=None, img_path=None):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if content:
        slide.placeholders[1].text = content
    if table:
        rows, cols = len(table), len(table[0])
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(3)
        tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for i in range(rows):
            for j in range(cols):
                tbl.cell(i,j).text = str(table[i][j])
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# --- Slide Content ---
slides = [
    ("SEI Executive KPI Report – High-Level Overview",
     "Prepared for: Executive Leadership\nPrepared by: SEI HR & Performance Analytics"),
    
    ("1. Industry Benchmark Intelligence",
     "Chargeability / Utilization Targets by Role:\n\nTechnical / Junior Consultants: 75–85%\nMid-Level Consultants / PMs: 75–90%\nSenior Leaders / Partners: 40–75%\nFirm Averages: 65–80%\n\nKey Insights:\n- Role-specific ranges increase transparency and employee buy-in.\n- Use utilization bands rather than rigid single targets.\n- Goldilocks Zone (2025): 70–80% chargeability for profitability and sustainability."),
    
    ("2. Broader KPIs Beyond Chargeability",
     "- Revenue per billable employee\n- Billable realization rate\n- Project margin\n- Client satisfaction / repeat business\n- Forecasted utilization / capacity planning\nRecommendation: Integrate KPIs like project margin or realization rate into leadership/PM scorecards."),
    
    ("3. Legal Compliance & Risk Mitigation",
     "- Protected Leaves: FMLA, CFRA, PDL, ADA/FEHA, PWFA, Jury duty, Voting leave, Bereavement, Military leave\n- Company Leaves: Vacation/PTO, Company holidays, Training/Professional Development\n- Non-Protected Leaves: Unpaid discretionary leave, Partial-day tardies, Non-statutory sick leave\nImplication: Automatically exclude protected leaves; apply pro-rated targets for extended leaves."),

    ("4. KPI Formula & Scoring",
     "Leave-Adjusted Chargeability Formula:\nAdjusted Chargeability (%) = Chargeable Hours / (Total Available Hours - Protected Leave Hours) * 100\n\nScoring Rubric:\n80%+ Exceeds Target\n65–79% Meets Target\n50–64% Needs Improvement\n<50% Does Not Meet Expectations\n\nBalanced Evaluation: 30–40% Chargeability + 60–70% Qualitative (Work Quality, Client Satisfaction, Project Delivery, Professional Development, Compliance & Teamwork)"),

    ("5. Fairness & Equity in KPI Design",
     "- Role-aligned expectations\n- Non-billable strategic work tracked but not penalized\n- Clear definitions: billable vs overhead\n- Regular review for work mix and seasonal variations\nEquity Audit: Quarterly/annual reviews, Disparate impact >10%, Documentation of justifications"),

    ("6. Time Tracking & Operational Best Practices",
     "- Clear, standardized task descriptions\n- Consistent time entry categories\n- Frequent audits\nWorkflow:\n1. Employees submit weekly timesheets\n2. HR flags protected/company leaves\n3. Adjusted Chargeability calculated\n4. KPI tables & summaries prepared\n5. Quarterly & annual executive review"),

    ("7. Visual Insights (Suggested Charts)",
     "Charts for Leadership:\n- KPI Bar Chart: Chargeability by employee/team\n- Gantt-like Timeline: High-level implementation plan\n- Heatmaps: Equity monitoring\n- Radar Charts: KPI balance across quality, client satisfaction, outcomes")
]

# Optional charts
if plt:
    # Example KPI Bar Chart
    kpi_chart_path = os.path.join(chart_folder, "kpi_bar_chart.png")
    plt.figure(figsize=(6,4))
    employees = ["Alice","Bob","Charlie","David"]
    chargeability = [85,78,90,70]
    plt.bar(employees, chargeability, color='skyblue')
    plt.ylim(0,100)
    plt.title("High-Level Chargeability by Employee")
    plt.savefig(kpi_chart_path)
    plt.close()
    slides.append(("KPI Visual", None, None, kpi_chart_path))

    # Example Gantt Timeline
    gantt_path = os.path.join(chart_folder, "gantt_timeline.png")
    tasks = ["Policy Design","Chart Generation","Manager Training","Employee Communication","Quarterly Audit"]
    start = [0,2,4,6,8]
    dur = [2,2,2,2,2]
    plt.figure(figsize=(6,3))
    for i, task in enumerate(tasks):
        plt.barh(task, dur[i], left=start[i], color='teal')
    plt.title("High-Level KPI Implementation Timeline")
    plt.savefig(gantt_path)
    plt.close()
    slides.append(("Implementation Timeline", None, None, gantt_path))

# Add slides
for s in slides:
    if len(s) == 2:
        add_slide(s[0], s[1])
    elif len(s) == 4:
        add_slide(s[0], s[1], s[2], s[3])

# Save presentation
prs.save(ppt_path)
print(f"✅ PowerPoint presentation generated at: {ppt_path}")
